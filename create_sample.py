#!/usr/bin/env python3
"""
Create a sample PowerPoint with intentional spelling/grammar errors for testing.
"""

from pptx import Presentation
from pptx.util import Inches, Pt

def main():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Quarterly Prezentation - Q4 Results"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Key Achievments:"

    p = tf.add_paragraph()
    p.text = "Recieved positive feedback from stakholders"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Thier team definately exceeded expectations"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Occured improvements in teh workflow"
    p.level = 1

    sample_file = "sample_with_errors.pptx"
    prs.save(sample_file)
    print(f"Created: {sample_file}")
    print("Run: python ppt_corrector.py sample_with_errors.pptx")

if __name__ == "__main__":
    main()
