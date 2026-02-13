#!/usr/bin/env python3
"""
End-to-end test: simulates the full Langdock workflow.

Steps:
  1. Create a sample .pptx with spelling errors
  2. Upload it to POST /api/upload (simulating the Langdock custom integration)
  3. Call MCP tool correct_pptx(file_id)
  4. Verify the download URL works (one-time download)
  5. Verify the download URL returns 404 on second access (one-time)
"""

import base64
import json
import os
import sys
import tempfile
from pathlib import Path

import httpx
from dotenv import load_dotenv
from pptx import Presentation

SERVER = "https://ppt-corrector-production.up.railway.app"

load_dotenv()

MCP_API_KEY = os.getenv("MCP_API_KEY", "").strip()
if not MCP_API_KEY:
    raise RuntimeError("Set MCP_API_KEY in environment or .env before running this test")

MCP_HEADERS = {
    "Content-Type": "application/json",
    "Accept": "application/json, text/event-stream",
    "Authorization": f"Bearer {MCP_API_KEY}",
}


def create_sample_pptx() -> bytes:
    """Create a small .pptx with intentional errors."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Quarterly Prezentation"
    slide.placeholders[1].text_frame.text = "Recieved feedback from stakholders"

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        prs.save(f.name)
        return Path(f.name).read_bytes()


def step_upload(file_bytes: bytes) -> str:
    """Step 1: Upload .pptx and get file_id."""
    print("=== Step 1: Upload .pptx ===")
    b64 = base64.b64encode(file_bytes).decode()
    r = httpx.post(
        f"{SERVER}/api/upload",
        json={"filename": "test_errors.pptx", "base64": b64},
        headers={
            "Authorization": f"Bearer {MCP_API_KEY}",
            "Content-Type": "application/json",
        },
        timeout=30,
    )
    assert r.status_code == 200, f"Upload failed: {r.status_code} {r.text}"
    data = r.json()
    file_id = data["file_id"]
    print(f"  file_id: {file_id}")
    print(f"  filename: {data['filename']}")
    print(f"  size: {data['size_bytes']} bytes")
    print(f"  expires in: {data['expires_in_minutes']} min")
    return file_id


def step_mcp_initialize() -> str:
    """Step 2: Initialize MCP session."""
    print("\n=== Step 2: MCP Initialize ===")
    r = httpx.post(
        f"{SERVER}/mcp",
        json={
            "jsonrpc": "2.0",
            "id": 1,
            "method": "initialize",
            "params": {
                "protocolVersion": "2025-03-26",
                "capabilities": {},
                "clientInfo": {"name": "e2e-test", "version": "1.0"},
            },
        },
        headers=MCP_HEADERS,
        timeout=15,
    )
    assert r.status_code == 200, f"Init failed: {r.status_code} {r.text}"
    sid = r.headers.get("mcp-session-id", "")
    print(f"  Session: {sid[:20]}..." if sid else "  Session: stateless")
    return sid


def step_call_tool(file_id: str, session_id: str) -> dict:
    """Step 3: Call correct_pptx MCP tool."""
    print("\n=== Step 3: Call correct_pptx ===")
    headers = {**MCP_HEADERS}
    if session_id:
        headers["mcp-session-id"] = session_id
    r = httpx.post(
        f"{SERVER}/mcp",
        json={
            "jsonrpc": "2.0",
            "id": 3,
            "method": "tools/call",
            "params": {
                "name": "correct_pptx",
                "arguments": {"file_id": file_id},
            },
        },
        headers=headers,
        timeout=300,  # processing can take minutes for large files
    )
    assert r.status_code == 200, f"Tool call failed: {r.status_code} {r.text}"

    # Parse SSE response
    result_data = None
    for line in r.text.split("\n"):
        if line.startswith("data:"):
            try:
                d = json.loads(line[5:].strip())
                if "result" in d:
                    result_data = d["result"]
            except json.JSONDecodeError:
                pass

    assert result_data is not None, f"No result in response: {r.text[:500]}"

    # The tool returns a JSON string inside content[0].text
    content = result_data.get("content", [])
    assert len(content) > 0, f"Empty content: {result_data}"
    tool_output = json.loads(content[0]["text"])

    print(f"  download_url: {tool_output['download_url']}")
    print(f"  filename: {tool_output['filename']}")
    print(f"  corrections: {tool_output['corrections_made']}")
    print(f"  slides: {tool_output['total_slides']}")
    if tool_output.get("changes"):
        for c in tool_output["changes"]:
            print(f"    Slide {c['slide']}: '{c['original']}' -> '{c['corrected']}'")

    return tool_output


def step_download(download_url: str) -> None:
    """Step 4: Download the corrected file (one-time)."""
    print("\n=== Step 4: Download (first time) ===")
    r = httpx.get(download_url, timeout=30)
    assert r.status_code == 200, f"Download failed: {r.status_code} {r.text}"
    assert len(r.content) > 100, f"Downloaded file too small: {len(r.content)} bytes"

    # Verify it's a valid .pptx (ZIP)
    assert r.content[:4] == b"PK\x03\x04", "Downloaded file is not a valid ZIP/PPTX"
    print(f"  Downloaded: {len(r.content)} bytes (valid .pptx)")

    # Check Content-Disposition header
    cd = r.headers.get("content-disposition", "")
    print(f"  Content-Disposition: {cd}")

    print("\n=== Step 5: Download again (should be 404) ===")
    r2 = httpx.get(download_url, timeout=10)
    assert r2.status_code == 404, f"Expected 404, got: {r2.status_code}"
    print(f"  Status: {r2.status_code} (one-time download confirmed)")


def main():
    print("Creating sample .pptx with errors...\n")
    file_bytes = create_sample_pptx()
    print(f"Sample file: {len(file_bytes)} bytes\n")

    file_id = step_upload(file_bytes)
    session_id = step_mcp_initialize()
    tool_output = step_call_tool(file_id, session_id)
    step_download(tool_output["download_url"])

    print("\n" + "=" * 50)
    print("ALL E2E TESTS PASSED!")
    print("=" * 50)


if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        print(f"\nTEST FAILED: {e}", file=sys.stderr)
        sys.exit(1)
