from __future__ import annotations

from io import BytesIO

from pptx import Presentation

from app.corrector import (
    apply_correction_to_runs,
    build_output_filename,
    correct_presentation_bytes,
    extract_paragraphs_from_shape,
    validate_pptx_bytes,
)
from tests.conftest import FakeLLMClient


def test_validate_pptx_bytes_accepts_valid_pptx(sample_pptx_bytes: bytes) -> None:
    safe_name = validate_pptx_bytes(sample_pptx_bytes, "../presentation.pptx")
    assert safe_name == "presentation.pptx"


def test_validate_pptx_bytes_rejects_corrupt_input() -> None:
    try:
        validate_pptx_bytes(b"not-a-pptx", "bad.pptx")
    except ValueError as exc:
        assert "valid .pptx" in str(exc) or "likely empty or corrupt" in str(exc)
    else:
        raise AssertionError("Expected corrupt input to be rejected")


def test_build_output_filename_uses_corrected_suffix() -> None:
    assert build_output_filename("folder/My Deck.pptx") == "My Deck_corrected.pptx"


def test_apply_correction_to_runs_preserves_formatting(sample_pptx_bytes: bytes) -> None:
    prs = Presentation(BytesIO(sample_pptx_bytes))
    paragraph = prs.slides[0].shapes[0].text_frame.paragraphs[0]

    original_bold = paragraph.runs[0].font.bold
    original_italic = paragraph.runs[1].font.italic
    original_size = paragraph.runs[2].font.size

    apply_correction_to_runs(paragraph, "The strategic objectives for Q3")

    assert paragraph.text == "The strategic objectives for Q3"
    assert paragraph.runs[0].font.bold == original_bold
    assert paragraph.runs[1].font.italic == original_italic
    assert paragraph.runs[2].font.size == original_size


def test_extract_paragraphs_from_shape_skips_shapes_without_text(sample_pptx_bytes: bytes) -> None:
    prs = Presentation(BytesIO(sample_pptx_bytes))
    slide = prs.slides[0]

    texts = []
    for shape in slide.shapes:
        for entry in extract_paragraphs_from_shape(shape):
            texts.append(entry["text"])

    assert "Teh strategc obiectives for Q3" in texts
    assert all(text for text in texts)


def test_correct_presentation_bytes_preserves_formatting_and_counts_changes(
    sample_pptx_bytes: bytes,
    fake_llm_client: FakeLLMClient,
) -> None:
    result = correct_presentation_bytes(
        sample_pptx_bytes,
        "proposal.pptx",
        fake_llm_client,
    )

    assert result.corrections_count == 2
    assert result.total_slides == 2
    assert result.file_name == "proposal_corrected.pptx"
    assert len(result.changes) == 2

    corrected_prs = Presentation(BytesIO(result.file_bytes))
    first_paragraph = corrected_prs.slides[0].shapes[0].text_frame.paragraphs[0]
    assert first_paragraph.text == "The strategic objectives for Q3"
    assert first_paragraph.runs[0].font.bold is True
    assert first_paragraph.runs[1].font.italic is True


def test_correct_presentation_bytes_returns_no_changes_for_identical_text(
    sample_pptx_bytes: bytes,
) -> None:
    llm_client = FakeLLMClient()

    result = correct_presentation_bytes(sample_pptx_bytes, "deck.pptx", llm_client)

    assert result.corrections_count == 0
    assert result.changes == []
    assert result.file_name == "deck_corrected.pptx"


def test_correct_presentation_bytes_handles_empty_presentation(empty_pptx_bytes: bytes) -> None:
    llm_client = FakeLLMClient()

    result = correct_presentation_bytes(empty_pptx_bytes, "empty.pptx", llm_client)

    assert result.corrections_count == 0
    assert result.total_slides == 0
    assert result.changes == []


def test_correct_presentation_bytes_handles_presentation_with_no_text(
    image_only_pptx_bytes: bytes,
) -> None:
    llm_client = FakeLLMClient()

    result = correct_presentation_bytes(image_only_pptx_bytes, "images-only.pptx", llm_client)

    assert result.corrections_count == 0
    assert result.total_slides == 1
    assert result.changes == []

