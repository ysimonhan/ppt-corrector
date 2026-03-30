from __future__ import annotations

import base64
from io import BytesIO
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from app.config import Settings
from app.runtime import AppRuntime, build_runtime
from app.storage import MemoryObjectStorage


def build_sample_pptx_bytes() -> bytes:
    prs = Presentation()

    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(7), Inches(1.5))
    paragraph = textbox.text_frame.paragraphs[0]
    run1 = paragraph.add_run()
    run1.text = "Teh "
    run1.font.bold = True
    run1.font.size = Pt(28)

    run2 = paragraph.add_run()
    run2.text = "strategc "
    run2.font.italic = True
    run2.font.size = Pt(20)

    run3 = paragraph.add_run()
    run3.text = "obiectives for Q3"
    run3.font.size = Pt(24)

    slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3), Inches(2), Inches(1))

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    table = slide2.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(5), Inches(1.5)).table
    cell_paragraph = table.cell(0, 0).text_frame.paragraphs[0]
    cell_run = cell_paragraph.add_run()
    cell_run.text = "Recieved feedback from stakholders"
    cell_run.font.size = Pt(22)

    slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(3), Inches(1), Inches(1))

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def build_empty_pptx_bytes() -> bytes:
    buffer = BytesIO()
    Presentation().save(buffer)
    return buffer.getvalue()


def build_image_only_pptx_bytes() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(3), Inches(2))
    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


class FakeLLMClient:
    def __init__(self, mapping: dict[str, str] | None = None, *, min_text_length: int = 3) -> None:
        self.mapping = mapping or {}
        self.min_text_length = min_text_length
        self.calls: list[str] = []

    def correct_text(self, text: str) -> tuple[str, bool]:
        self.calls.append(text)
        if text in self.mapping:
            return self.mapping[text], True
        return text, True


class RecordingJobQueue:
    def __init__(self) -> None:
        self.job_ids: list[str] = []

    def enqueue(self, job_id: str) -> None:
        self.job_ids.append(job_id)


@pytest.fixture()
def sample_pptx_bytes() -> bytes:
    return build_sample_pptx_bytes()


@pytest.fixture()
def empty_pptx_bytes() -> bytes:
    return build_empty_pptx_bytes()


@pytest.fixture()
def image_only_pptx_bytes() -> bytes:
    return build_image_only_pptx_bytes()


@pytest.fixture()
def sample_pptx_base64(sample_pptx_bytes: bytes) -> str:
    return base64.b64encode(sample_pptx_bytes).decode("utf-8")


@pytest.fixture()
def test_settings(tmp_path: Path) -> Settings:
    db_path = tmp_path / "ppt_corrector_test.db"
    return Settings(
        api_key="test-api-key",
        langdock_api_key="test-langdock-key",
        langdock_api_url="https://example.com/langdock",
        langdock_model="gpt-5.4",
        port=8000,
        min_text_length=3,
        llm_timeout_seconds=1.0,
        max_upload_size_bytes=10 * 1024 * 1024,
        job_ttl_seconds=600,
        job_cleanup_interval_seconds=60,
        metadata_retention_seconds=3600,
        default_highlight_color="FFFF00",
        database_url=f"sqlite:///{db_path}",
        queue_backend="inline",
        redis_url="redis://localhost:6379/0",
        rq_queue_name="ppt-corrector-test",
        storage_backend="memory",
        s3_bucket="",
        s3_region="eu-central-1",
        s3_endpoint_url="",
        s3_access_key_id="",
        s3_secret_access_key="",
        s3_prefix="ppt-corrector-test",
    )


@pytest.fixture()
def fake_llm_client() -> FakeLLMClient:
    return FakeLLMClient(
        {
            "Teh strategc obiectives for Q3": "The strategic objectives for Q3",
            "Recieved feedback from stakholders": "Received feedback from stakeholders",
        }
    )


@pytest.fixture()
def recording_queue() -> RecordingJobQueue:
    return RecordingJobQueue()


@pytest.fixture()
def memory_storage() -> MemoryObjectStorage:
    return MemoryObjectStorage()


@pytest.fixture()
def runtime(
    test_settings: Settings,
    recording_queue: RecordingJobQueue,
    memory_storage: MemoryObjectStorage,
    fake_llm_client: FakeLLMClient,
) -> AppRuntime:
    return build_runtime(
        test_settings,
        storage=memory_storage,
        queue=recording_queue,
        llm_client_factory=lambda: fake_llm_client,
    )
