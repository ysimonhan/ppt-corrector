from __future__ import annotations

from io import BytesIO

from lxml import etree
from pptx import Presentation

from app.highlighter import apply_correction_to_runs_highlighted, diff_segments, normalize_highlight_color


def test_normalize_highlight_color() -> None:
    assert normalize_highlight_color("#ff9632") == "FF9632"


def test_normalize_highlight_color_rejects_invalid_value() -> None:
    try:
        normalize_highlight_color("not-a-color")
    except ValueError as exc:
        assert "hex RGB" in str(exc)
    else:
        raise AssertionError("Expected invalid color to raise ValueError")


def test_diff_segments_marks_changed_words() -> None:
    segments = diff_segments("Recieved feedback from stakholders", "Received feedback from stakeholders")

    assert any(changed for _, changed in segments)
    assert "".join(text for text, _ in segments) == "Received feedback from stakeholders"


def test_apply_correction_to_runs_highlighted_injects_xml_highlight(sample_pptx_bytes: bytes) -> None:
    prs = Presentation(BytesIO(sample_pptx_bytes))
    paragraph = prs.slides[0].shapes[0].text_frame.paragraphs[0]

    apply_correction_to_runs_highlighted(
        paragraph,
        "The strategic objectives for Q3",
        "FFFF00",
    )

    assert paragraph.text == "The strategic objectives for Q3"

    xml = etree.tostring(paragraph._p, encoding="unicode")
    assert "<a:highlight>" in xml
    assert "FFFF00" in xml

