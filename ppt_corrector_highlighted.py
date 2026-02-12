#!/usr/bin/env python3
"""
PowerPoint Spelling & Grammar Correction Agent – Highlighted Version

Same as ppt_corrector.py but additionally highlights corrected text with a
yellow marker so reviewers can easily spot what was changed.

Usage:
    python ppt_corrector_highlighted.py presentation.pptx
    python ppt_corrector_highlighted.py presentation.pptx --output custom_name.pptx
    python ppt_corrector_highlighted.py presentation.pptx --color FF9632  # orange marker
"""

import argparse
import copy
import difflib
import json
import logging
import os
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import httpx
from dotenv import load_dotenv
from lxml.etree import SubElement
from pptx import Presentation
from tenacity import retry, stop_after_attempt, wait_exponential

# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
)
logger = logging.getLogger(__name__)

# Langdock API - uses Bearer token per https://docs.langdock.com/api-endpoints/api-introduction
LANGDOCK_URL = "https://api.langdock.com/anthropic/eu/v1/messages"
LANGDOCK_MODEL = "claude-sonnet-4-5-20250929"
MIN_TEXT_LENGTH = 3
DEFAULT_HIGHLIGHT_COLOR = "FFFF00"  # Yellow marker

# DrawingML namespace used by PowerPoint XML
_nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


def get_api_key() -> str:
    """Get and validate Langdock API key."""
    api_key = os.getenv("LANGDOCK_API_KEY")
    if not api_key:
        raise ValueError(
            "LANGDOCK_API_KEY not set. Add it to .env or export it. "
            "Get your key from https://app.langdock.com"
        )
    key = api_key.strip()
    if not key.startswith("sk-"):
        logger.warning("API key should usually start with 'sk-'. Check you copied the full key.")
    return key


SYSTEM_PROMPT = """You are a highly rated and experienced Engagement Manager from McKinsey & Company in Germany. You are super good at correcting spelling and grammar errors in texts.

<Task>
Your task is to correct ONLY spelling and grammar errors in the given text.
</Task>

<Context>
- You will often only get one word for correction. This is normal and expected since the data comes from a powerpoint presentation where the text needs to be extraced from each
object individually so often there is only one word. If there is no text, return an empty string. If you only get one word, return the corrected word.

Rules:
- Return ONLY the corrected text, nothing else. No explanations, no quotes, no preamble. 
- As context one word is enough. Do not ask questions or make assumptions, just correct the word.
- Preserve technical terms, proper nouns, brand names, and acronyms.
- Do not change formatting, punctuation style, or sentence structure unless grammatically wrong.
- If the text has no errors, return it unchanged.
- Keep the output the same length and style as the input.

"""


@retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=1, min=2, max=10))
def correct_text_with_llm(text: str, api_key: str) -> Tuple[str, bool]:
    """
    Correct spelling and grammar using Claude Sonnet 4.5 via Langdock API.
    Uses direct HTTP with Bearer token (Langdock requires Authorization: Bearer).

    Returns:
        Tuple of (corrected_text, success_flag)
    """
    text = text.strip()
    if len(text) < MIN_TEXT_LENGTH:
        return text, True

    try:
        payload = {
            "model": LANGDOCK_MODEL,
            "max_tokens": 512,
            "temperature": 0.2,
            "system": SYSTEM_PROMPT,
            "messages": [
                {
                    "role": "user",
                    "content": f"Correct spelling and grammar in this text. Return ONLY the corrected text:\n\n{text}",
                }
            ],
        }
        with httpx.Client(timeout=60.0) as client:
            resp = client.post(
                LANGDOCK_URL,
                headers={
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json",
                },
                json=payload,
            )

        if resp.status_code == 401:
            logger.error(
                "401 Unauthorized: API key invalid. Check:\n"
                "  1. Key is correct in .env (Settings -> API in Langdock)\n"
                "  2. Key has scope for Completion/Anthropic API\n"
                "  3. Create a new key at https://app.langdock.com if needed"
            )
            return text, False

        resp.raise_for_status()
        data = resp.json()

        corrected = ""
        for block in data.get("content", []):
            if block.get("type") == "text":
                corrected += block.get("text", "") or ""

        corrected = corrected.strip()
        if not corrected:
            logger.warning(f"Empty response for: {text[:50]}...")
            return text, True

        return corrected, True

    except httpx.HTTPStatusError as e:
        logger.error(f"LLM API error: {e.response.status_code} - {e.response.text[:200]}")
        return text, False
    except Exception as e:
        logger.error(f"LLM error for '{text[:30]}...': {e}")
        return text, False


# ---------------------------------------------------------------------------
#  Highlighting helpers – word-level diff
# ---------------------------------------------------------------------------

_A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _highlight_run_xml(r_elem: Any, color_hex: str) -> None:
    """Apply highlight marker to a raw <a:r> lxml element."""
    rPr = r_elem.find(f"{_A_NS}rPr")
    if rPr is None:
        rPr = SubElement(r_elem, f"{_A_NS}rPr")
        # rPr must come before <a:t>, reorder
        r_elem.remove(rPr)
        r_elem.insert(0, rPr)

    # Remove existing highlight if present
    for existing in rPr.findall(f"{_A_NS}highlight"):
        rPr.remove(existing)

    hl = SubElement(rPr, f"{_A_NS}highlight")
    srgb = SubElement(hl, f"{_A_NS}srgbClr")
    srgb.set("val", color_hex)


def _clone_run(r_elem: Any, new_text: str) -> Any:
    """Deep-copy a run element and set its text, preserving all formatting."""
    new_r = copy.deepcopy(r_elem)
    t_elem = new_r.find(f"{_A_NS}t")
    if t_elem is not None:
        t_elem.text = new_text
    return new_r


def _diff_segments(original: str, corrected: str) -> List[Tuple[str, bool]]:
    """Word-level diff returning segments with a changed flag.

    Returns list of (text, is_changed) tuples that concatenate to corrected.
    Unchanged words keep is_changed=False; replaced/inserted words get True.
    Spaces between words are attached to the preceding segment.

    Example:
        original:  "Recieved positive feedback from stakholders"
        corrected: "Received positive feedback from stakeholders"
        -> [("Received ", True), ("positive feedback from ", False), ("stakeholders", True)]
    """
    orig_words = original.split()
    corr_words = corrected.split()

    sm = difflib.SequenceMatcher(None, orig_words, corr_words)
    segments: List[Tuple[str, bool]] = []

    for op, _i1, _i2, j1, j2 in sm.get_opcodes():
        if op == "equal":
            text = " ".join(corr_words[j1:j2])
            segments.append((text, False))
        else:
            # replace, insert, or delete (on corrected side)
            if j1 < j2:
                text = " ".join(corr_words[j1:j2])
                segments.append((text, True))

    # Re-join with spaces between segments
    result: List[Tuple[str, bool]] = []
    for idx, (text, changed) in enumerate(segments):
        if idx < len(segments) - 1:
            text += " "
        result.append((text, changed))

    return result


def apply_correction_to_runs_highlighted(
    paragraph: Any, corrected_text: str, color_hex: str
) -> None:
    """Apply corrected text back to runs, highlighting ONLY changed words.

    1. Collect original text from all runs.
    2. Diff original vs corrected at word level.
    3. For each diff segment, create a run cloned from the first original run
       (preserving font/size/color) and apply highlight only to changed segments.
    4. Replace all original runs in the paragraph XML with the new split runs.
    """
    runs = list(paragraph.runs)
    if not runs:
        return

    original_text = "".join(r.text for r in runs)

    # If nothing actually changed after stripping, bail out
    if original_text.strip() == corrected_text.strip():
        return

    segments = _diff_segments(original_text, corrected_text)

    # Use the first run as the formatting template
    template_r = runs[0]._r

    # Build new run elements
    new_run_elems = []
    for seg_text, is_changed in segments:
        if not seg_text:
            continue
        new_r = _clone_run(template_r, seg_text)
        if is_changed:
            _highlight_run_xml(new_r, color_hex)
        new_run_elems.append(new_r)

    # Replace original runs in the paragraph XML
    p_elem = paragraph._p

    # Remove old <a:r> elements
    for run in runs:
        p_elem.remove(run._r)

    # Append new runs (after any existing elements like <a:pPr>)
    for new_r in new_run_elems:
        p_elem.append(new_r)


# ---------------------------------------------------------------------------
#  Shape / paragraph extraction (same as ppt_corrector.py)
# ---------------------------------------------------------------------------

def extract_paragraphs_from_shape(shape: Any) -> List[Dict[str, Any]]:
    """Extract paragraphs from a PowerPoint shape (text frames, tables, grouped shapes)."""
    paragraphs: List[Dict[str, Any]] = []

    try:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                full_text = paragraph.text.strip()
                if full_text:
                    paragraphs.append({"text": full_text, "paragraph": paragraph})

        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        full_text = paragraph.text.strip()
                        if full_text:
                            paragraphs.append({"text": full_text, "paragraph": paragraph})

        # Grouped shapes (MSO_SHAPE_TYPE.GROUP = 6)
        if hasattr(shape, "shape_type") and shape.shape_type == 6:
            for sub_shape in shape.shapes:
                paragraphs.extend(extract_paragraphs_from_shape(sub_shape))

    except Exception as e:
        logger.warning(f"Error extracting from shape: {e}")

    return paragraphs


# ---------------------------------------------------------------------------
#  Main correction function
# ---------------------------------------------------------------------------

def correct_powerpoint_highlighted(
    input_file: str,
    output_file: Optional[str] = None,
    save_report: bool = True,
    highlight_color: str = DEFAULT_HIGHLIGHT_COLOR,
) -> Dict[str, Any]:
    """
    Correct PowerPoint spelling/grammar and highlight every change with a marker.

    Args:
        input_file:       Path to input PowerPoint file
        output_file:      Path to output file (default: {stem}_v_correctedbyai.pptx)
        save_report:      Whether to save a JSON report of changes
        highlight_color:  Hex RGB color for the marker (default: FFFF00 = yellow)

    Returns:
        Dict with status, output_file, statistics, and changes
    """
    input_path = Path(input_file)

    if not input_path.exists():
        raise FileNotFoundError(f"File not found: {input_file}")

    if input_path.suffix.lower() not in (".pptx", ".ppt"):
        raise ValueError(f"Invalid file type: {input_path.suffix}. Use .pptx")

    if output_file is None:
        output_file = str(
            input_path.parent / f"{input_path.stem}_v_correctedbyai{input_path.suffix}"
        )

    logger.info(f"Loading presentation: {input_file}")
    prs = Presentation(input_file)
    api_key = get_api_key()

    corrections_made = 0
    total_slides = len(prs.slides)
    change_log: List[Dict[str, Any]] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        logger.info(f"Processing slide {slide_num}/{total_slides}")

        for shape_idx, shape in enumerate(slide.shapes):
            paragraphs = extract_paragraphs_from_shape(shape)

            for element in paragraphs:
                original_text = element["text"]

                if len(original_text.strip()) < MIN_TEXT_LENGTH:
                    continue

                corrected_text, success = correct_text_with_llm(original_text, api_key)

                if not success:
                    continue

                if corrected_text.strip() != original_text.strip():
                    # Write corrected text into runs AND apply highlight marker
                    apply_correction_to_runs_highlighted(
                        element["paragraph"], corrected_text, highlight_color
                    )
                    corrections_made += 1
                    change_log.append({
                        "slide": slide_num,
                        "shape": shape_idx,
                        "original": original_text,
                        "corrected": corrected_text,
                    })
                    logger.info(f"  ✓ Slide {slide_num}: '{original_text}' -> '{corrected_text}'")

    logger.info(f"Saving: {output_file}")
    prs.save(output_file)

    result = {
        "status": "success",
        "output_file": output_file,
        "corrections_made": corrections_made,
        "total_slides": total_slides,
        "highlight_color": f"#{highlight_color}",
        "changes": change_log,
    }

    if save_report:
        report_path = Path(output_file).parent / f"{Path(output_file).stem}_report.json"
        with open(report_path, "w", encoding="utf-8") as f:
            json.dump(result, f, indent=2, ensure_ascii=False)
        logger.info(f"Report saved: {report_path}")

    logger.info(f"Done! {corrections_made} corrections highlighted across {total_slides} slides")
    return result


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Correct PowerPoint spelling/grammar and highlight changes with a yellow marker"
    )
    parser.add_argument(
        "input_file",
        help="Path to PowerPoint file (.pptx)",
    )
    parser.add_argument(
        "-o", "--output",
        help="Output file path (default: {input}_v_correctedbyai.pptx)",
    )
    parser.add_argument(
        "--no-report",
        action="store_true",
        help="Skip saving the JSON correction report",
    )
    parser.add_argument(
        "--color",
        default=DEFAULT_HIGHLIGHT_COLOR,
        help="Highlight color as hex RGB, e.g. FFFF00 for yellow (default: FFFF00)",
    )
    parser.add_argument(
        "-v", "--verbose",
        action="store_true",
        help="Enable debug logging",
    )

    args = parser.parse_args()

    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)

    try:
        result = correct_powerpoint_highlighted(
            args.input_file,
            output_file=args.output,
            save_report=not args.no_report,
            highlight_color=args.color.lstrip("#"),
        )
        print(f"\nCorrected file saved to: {result['output_file']}")
    except Exception as e:
        logger.error(f"Error: {e}")
        raise SystemExit(1)


if __name__ == "__main__":
    main()
