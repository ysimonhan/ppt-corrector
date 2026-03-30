from __future__ import annotations

import logging
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Any

from pptx import Presentation

from app.highlighter import DEFAULT_HIGHLIGHT_COLOR, apply_correction_to_runs_highlighted, normalize_highlight_color
from app.llm import LangdockLLMClient


logger = logging.getLogger(__name__)


@dataclass(frozen=True)
class PresentationCorrectionResult:
    file_bytes: bytes
    file_name: str
    corrections_count: int
    total_slides: int
    changes: list[dict[str, Any]]


def validate_pptx_bytes(file_bytes: bytes, file_name: str) -> str:
    safe_name = Path(file_name).name or "presentation.pptx"
    if len(file_bytes) < 100:
        raise ValueError(f"File is only {len(file_bytes)} bytes - likely empty or corrupt.")
    if file_bytes[:4] != b"PK\x03\x04":
        raise ValueError("File is not a valid .pptx (ZIP) package.")
    if not safe_name.lower().endswith(".pptx"):
        raise ValueError(f"Unsupported extension: {safe_name}. Only .pptx is accepted.")
    return safe_name


def build_output_filename(file_name: str) -> str:
    safe_name = Path(file_name).name or "presentation.pptx"
    stem = Path(safe_name).stem or "presentation"
    return f"{stem}_corrected.pptx"


def apply_correction_to_runs(paragraph: Any, corrected_text: str) -> None:
    """Apply corrected text back to paragraph runs, preserving each run's formatting."""
    runs = list(paragraph.runs)

    if not runs:
        return

    if len(runs) == 1:
        runs[0].text = corrected_text
        return

    original_lengths = [len(run.text) for run in runs]
    total_original = sum(original_lengths)

    if total_original == 0:
        runs[0].text = corrected_text
        return

    position = 0
    for index, run in enumerate(runs):
        if index == len(runs) - 1:
            run.text = corrected_text[position:]
        else:
            share = round(len(corrected_text) * original_lengths[index] / total_original)
            run.text = corrected_text[position : position + share]
            position += share


def extract_paragraphs_from_shape(shape: Any) -> list[dict[str, Any]]:
    """Extract paragraphs from a PowerPoint shape."""
    paragraphs: list[dict[str, Any]] = []

    try:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                full_text = paragraph.text.strip()
                if full_text:
                    paragraphs.append({"text": full_text, "paragraph": paragraph})

        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        full_text = paragraph.text.strip()
                        if full_text:
                            paragraphs.append({"text": full_text, "paragraph": paragraph})

        if hasattr(shape, "shape_type") and shape.shape_type == 6:
            for sub_shape in shape.shapes:
                paragraphs.extend(extract_paragraphs_from_shape(sub_shape))

    except Exception as exc:
        logger.warning("Error extracting from shape: %s", exc)

    return paragraphs


def correct_presentation_bytes(
    file_bytes: bytes,
    file_name: str,
    llm_client: LangdockLLMClient,
    *,
    highlight: bool = False,
    highlight_color: str = DEFAULT_HIGHLIGHT_COLOR,
) -> PresentationCorrectionResult:
    safe_name = validate_pptx_bytes(file_bytes, file_name)
    normalized_color = normalize_highlight_color(highlight_color)

    presentation = Presentation(BytesIO(file_bytes))

    corrections_made = 0
    total_slides = len(presentation.slides)
    change_log: list[dict[str, Any]] = []

    for slide_number, slide in enumerate(presentation.slides, 1):
        logger.info("Processing slide %s/%s", slide_number, total_slides)

        for shape in slide.shapes:
            paragraphs = extract_paragraphs_from_shape(shape)

            for element in paragraphs:
                original_text = element["text"]
                if len(original_text.strip()) < llm_client.min_text_length:
                    continue

                corrected_text, success = llm_client.correct_text(original_text)
                if not success:
                    continue

                if corrected_text.strip() != original_text.strip():
                    if highlight:
                        apply_correction_to_runs_highlighted(
                            element["paragraph"],
                            corrected_text,
                            normalized_color,
                        )
                    else:
                        apply_correction_to_runs(element["paragraph"], corrected_text)

                    corrections_made += 1
                    change_log.append(
                        {
                            "slide": slide_number,
                            "original": original_text,
                            "corrected": corrected_text,
                        }
                    )

    output_buffer = BytesIO()
    presentation.save(output_buffer)

    return PresentationCorrectionResult(
        file_bytes=output_buffer.getvalue(),
        file_name=build_output_filename(safe_name),
        corrections_count=corrections_made,
        total_slides=total_slides,
        changes=change_log,
    )
