#!/usr/bin/env python3
"""
PowerPoint Spelling & Grammar Correction Agent

Reads a PowerPoint presentation, corrects spelling and grammar using Claude Sonnet 4.5
via Langdock API, and saves a "_v_correctedbyai" version in the same folder.

Usage:
    python ppt_corrector.py presentation.pptx
    python ppt_corrector.py presentation.pptx --output custom_name.pptx
"""

import argparse
import json
import logging
import os
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import httpx
from dotenv import load_dotenv
from pptx import Presentation
from tenacity import retry, stop_after_attempt, wait_exponential

# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
)
logger = logging.getLogger(__name__)

# Langdock API - uses Bearer token per https://docs.langdock.com/api-endpoints/api-introduction
LANGDOCK_URL = "https://api.langdock.com/anthropic/eu/v1/messages"
LANGDOCK_MODEL = "claude-sonnet-4-5-20250929"
MIN_TEXT_LENGTH = 3


def get_api_key() -> str:
    """Get and validate Langdock API key."""
    api_key = os.getenv("LANGDOCK_API_KEY")
    if not api_key:
        raise ValueError(
            "LANGDOCK_API_KEY not set. Add it to .env or export it. "
            "Get your key from https://app.langdock.com"
        )
    key = api_key.strip()
    if not key.startswith("sk-"):
        logger.warning("API key should usually start with 'sk-'. Check you copied the full key.")
    return key


SYSTEM_PROMPT = """You are a highly rated and experienced Engagement Manager from McKinsey & Company in Germany. You are super good at correcting spelling and grammar errors in texts.

<Task>
Your task is to correct ONLY spelling and grammar errors in the given text.
</Task>

<Context>
- You will often only get one word for correction. This is normal and expected since the data comes from a powerpoint presentation where the text needs to be extraced from each
object individually so often there is only one word. If there is no text, return an empty string. If you only get one word, return the corrected word.

Rules:
- Return ONLY the corrected text, nothing else. No explanations, no quotes, no preamble. 
- As context one word is enough. Do not ask questions or make assumptions, just correct the word.
- Preserve technical terms, proper nouns, brand names, and acronyms.
- Do not change formatting, punctuation style, or sentence structure unless grammatically wrong.
- If the text has no errors, return it unchanged.
- Keep the output the same length and style as the input.

"""


@retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=1, min=2, max=10))
def correct_text_with_llm(text: str, api_key: str) -> Tuple[str, bool]:
    """
    Correct spelling and grammar using Claude Sonnet 4.5 via Langdock API.
    Uses direct HTTP with Bearer token (Langdock requires Authorization: Bearer).

    Returns:
        Tuple of (corrected_text, success_flag)
    """
    text = text.strip()
    if len(text) < MIN_TEXT_LENGTH:
        return text, True

    try:
        # Langdock requires Authorization: Bearer - use direct HTTP for exact control
        payload = {
            "model": LANGDOCK_MODEL,
            "max_tokens": 512,
            "temperature": 0.2,
            "system": SYSTEM_PROMPT,
            "messages": [
                {
                    "role": "user",
                    "content": f"Correct spelling and grammar in this text. Return ONLY the corrected text:\n\n{text}",
                }
            ],
        }
        with httpx.Client(timeout=60.0) as client:
            resp = client.post(
                LANGDOCK_URL,
                headers={
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json",
                },
                json=payload,
            )

        if resp.status_code == 401:
            logger.error(
                "401 Unauthorized: API key invalid. Check:\n"
                "  1. Key is correct in .env (Settings → API in Langdock)\n"
                "  2. Key has scope for Completion/Anthropic API\n"
                "  3. Create a new key at https://app.langdock.com if needed"
            )
            return text, False

        resp.raise_for_status()
        data = resp.json()

        corrected = ""
        for block in data.get("content", []):
            if block.get("type") == "text":
                corrected += block.get("text", "") or ""

        corrected = corrected.strip()
        if not corrected:
            logger.warning(f"Empty response for: {text[:50]}...")
            return text, True

        return corrected, True

    except httpx.HTTPStatusError as e:
        logger.error(f"LLM API error: {e.response.status_code} - {e.response.text[:200]}")
        return text, False
    except Exception as e:
        logger.error(f"LLM error for '{text[:30]}...': {e}")
        return text, False


def apply_correction_to_runs(paragraph: Any, corrected_text: str) -> None:
    """Apply corrected text back to paragraph runs, preserving each run's formatting.

    Setting paragraph.text destroys all formatting. Instead we write corrected text
    back into the existing runs so font, size, color, bold, italic, etc. survive.

    Strategy:
    - 1 run  -> replace its text directly (perfect preservation).
    - N runs -> distribute corrected text proportionally across runs so each run
                keeps roughly the same share it had before. Every run's font
                properties stay intact because we only touch run.text.
    """
    runs = list(paragraph.runs)

    if not runs:
        return

    if len(runs) == 1:
        runs[0].text = corrected_text
        return

    # Multiple runs: distribute proportionally by original character count
    orig_lengths = [len(r.text) for r in runs]
    total_orig = sum(orig_lengths)

    if total_orig == 0:
        runs[0].text = corrected_text
        return

    pos = 0
    for i, run in enumerate(runs):
        if i == len(runs) - 1:
            # Last run gets whatever remains
            run.text = corrected_text[pos:]
        else:
            share = round(len(corrected_text) * orig_lengths[i] / total_orig)
            run.text = corrected_text[pos : pos + share]
            pos += share


def extract_paragraphs_from_shape(shape: Any) -> List[Dict[str, Any]]:
    """
    Extract paragraphs from a PowerPoint shape (text frames, tables, grouped shapes).

    Returns list of dicts with 'text' and 'paragraph' for each non-empty paragraph.
    """
    paragraphs: List[Dict[str, Any]] = []

    try:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                full_text = paragraph.text.strip()
                if full_text:
                    paragraphs.append({"text": full_text, "paragraph": paragraph})

        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        full_text = paragraph.text.strip()
                        if full_text:
                            paragraphs.append({"text": full_text, "paragraph": paragraph})

        # Grouped shapes (MSO_SHAPE_TYPE.GROUP = 6)
        if hasattr(shape, "shape_type") and shape.shape_type == 6:
            for sub_shape in shape.shapes:
                paragraphs.extend(extract_paragraphs_from_shape(sub_shape))

    except Exception as e:
        logger.warning(f"Error extracting from shape: {e}")

    return paragraphs


def correct_powerpoint(
    input_file: str,
    output_file: Optional[str] = None,
    save_report: bool = True,
) -> Dict[str, Any]:
    """
    Correct PowerPoint spelling and grammar using Claude Sonnet 4.5 via Langdock.

    Args:
        input_file: Path to input PowerPoint file
        output_file: Path to output file (default: {stem}_v_correctedbyai.pptx)
        save_report: Whether to save a JSON report of changes

    Returns:
        Dict with status, output_file, statistics, and changes
    """
    input_path = Path(input_file)

    if not input_path.exists():
        raise FileNotFoundError(f"File not found: {input_file}")

    if input_path.suffix.lower() not in (".pptx", ".ppt"):
        raise ValueError(f"Invalid file type: {input_path.suffix}. Use .pptx")

    if output_file is None:
        output_file = str(
            input_path.parent / f"{input_path.stem}_v_correctedbyai{input_path.suffix}"
        )

    logger.info(f"Loading presentation: {input_file}")
    prs = Presentation(input_file)
    api_key = get_api_key()

    corrections_made = 0
    total_slides = len(prs.slides)
    change_log: List[Dict[str, Any]] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        logger.info(f"Processing slide {slide_num}/{total_slides}")

        for shape_idx, shape in enumerate(slide.shapes):
            paragraphs = extract_paragraphs_from_shape(shape)

            for element in paragraphs:
                original_text = element["text"]

                if len(original_text.strip()) < MIN_TEXT_LENGTH:
                    continue

                corrected_text, success = correct_text_with_llm(original_text, api_key)

                if not success:
                    continue

                if corrected_text.strip() != original_text.strip():
                    # Write corrected text back into runs (preserves formatting)
                    apply_correction_to_runs(element["paragraph"], corrected_text)
                    corrections_made += 1
                    change_log.append({
                        "slide": slide_num,
                        "shape": shape_idx,
                        "original": original_text,
                        "corrected": corrected_text,
                    })
                    logger.info(f"  ✓ Slide {slide_num}: '{original_text}' → '{corrected_text}'")

    logger.info(f"Saving: {output_file}")
    prs.save(output_file)

    result = {
        "status": "success",
        "output_file": output_file,
        "corrections_made": corrections_made,
        "total_slides": total_slides,
        "changes": change_log,
    }

    if save_report:
        report_path = Path(output_file).parent / f"{Path(output_file).stem}_report.json"
        with open(report_path, "w", encoding="utf-8") as f:
            json.dump(result, f, indent=2, ensure_ascii=False)
        logger.info(f"Report saved: {report_path}")

    logger.info(f"✓ Complete! Made {corrections_made} corrections across {total_slides} slides")
    return result


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Correct PowerPoint spelling and grammar using Claude Sonnet 4.5 via Langdock API"
    )
    parser.add_argument(
        "input_file",
        help="Path to PowerPoint file (.pptx)",
    )
    parser.add_argument(
        "-o", "--output",
        help="Output file path (default: {input}_v_correctedbyai.pptx)",
    )
    parser.add_argument(
        "--no-report",
        action="store_true",
        help="Skip saving the JSON correction report",
    )
    parser.add_argument(
        "-v", "--verbose",
        action="store_true",
        help="Enable debug logging",
    )

    args = parser.parse_args()

    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)

    try:
        result = correct_powerpoint(
            args.input_file,
            output_file=args.output,
            save_report=not args.no_report,
        )
        print(f"\nCorrected file saved to: {result['output_file']}")
    except Exception as e:
        logger.error(f"Error: {e}")
        raise SystemExit(1)


if __name__ == "__main__":
    main()
